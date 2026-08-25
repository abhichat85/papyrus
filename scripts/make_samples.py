"""Build the demo documents served by the landing page.

These are deliberately *realistic* rather than minimal: a multi-page PDF
with a real table and a running header, a Word document with nested lists,
a deck with speaker notes, and a workbook with number formats. The demo is
only honest if the samples exercise the parts of the engine the page
claims work.
"""

from __future__ import annotations

from pathlib import Path

OUT = Path(__file__).resolve().parent.parent / "web" / "public" / "samples"

ROWS = [
    ["Metric", "FY 2023", "FY 2024", "FY 2025"],
    ["Revenue", "6,200,000", "10,040,000", "14,120,000"],
    ["Gross margin", "71%", "74%", "78%"],
    ["EBITDA", "310,000", "905,000", "2,140,000"],
    ["Net retention", "104%", "112%", "119%"],
    ["Headcount", "38", "61", "94"],
]

SUMMARY = [
    "Revenue grew forty one percent year over year, ahead of the conservative",
    "plan set at the start of the period. Growth came disproportionately from",
    "the enterprise segment, where net retention reached one hundred and",
    "nineteen percent on the back of two large expansions in the fourth",
    "quarter.",
]

OUTLOOK = [
    "The coming year compounds on this base. Three commitments shape the plan:",
    "hold gross margin above seventy five percent while absorbing the cost of",
    "the new region, keep headcount growth below revenue growth, and convert",
    "the pipeline built in the second half without discounting.",
]

RISKS = [
    "Concentration: the top ten accounts are thirty one percent of revenue.",
    "Currency: roughly a fifth of billings are non-dollar and unhedged.",
    "Hiring: two of the four regional leads are still open.",
]


def annual_report_pdf() -> None:
    import pymupdf

    doc = pymupdf.open()
    doc.set_metadata(
        {
            "title": "Annual Report 2025",
            "author": "Office of the CFO",
            "subject": "Audited results for the year ended 31 December 2025",
        }
    )

    def page_frame(page, number: int) -> None:
        # A running header and a page number on every page — exactly the
        # furniture Papyrus is supposed to notice and drop.
        page.insert_text((72, 46), "Annual Report 2025  |  Confidential", fontsize=8, color=(0.4, 0.4, 0.4))
        page.draw_line(pymupdf.Point(72, 54), pymupdf.Point(523, 54), color=(0.8, 0.8, 0.8), width=0.5)
        page.insert_text((295, 792), str(number), fontsize=8, color=(0.4, 0.4, 0.4))

    page = doc.new_page()
    page_frame(page, 1)
    page.insert_text((72, 120), "Annual Report", fontsize=30, fontname="hebo")
    page.insert_text((72, 152), "Financial year 2025", fontsize=14, color=(0.35, 0.35, 0.35))
    y = 210
    page.insert_text((72, y), "Executive Summary", fontsize=17, fontname="hebo")
    y += 28
    for line in SUMMARY:
        page.insert_text((72, y), line, fontsize=11)
        y += 17
    y += 18
    page.insert_text((72, y), "Highlights", fontsize=13, fontname="hebo")
    y += 24
    for bullet in (
        "- Enterprise pipeline doubled against the prior year",
        "- Gross churn fell to one point two percent",
        "- Two new regions opened, both ahead of schedule",
        "- First full year of positive operating cash flow",
    ):
        page.insert_text((80, y), bullet, fontsize=11)
        y += 17

    page = doc.new_page()
    page_frame(page, 2)
    y = 110
    page.insert_text((72, y), "Financial Performance", fontsize=17, fontname="hebo")
    y += 32
    left, widths, row_height = 72, [150, 100, 100, 100], 26
    for index, row in enumerate(ROWS):
        x = left
        top = y + index * row_height
        for width, cell in zip(widths, row, strict=False):
            rect = pymupdf.Rect(x, top, x + width, top + row_height)
            page.draw_rect(rect, color=(0.75, 0.75, 0.75), width=0.6)
            page.insert_text(
                (x + 8, top + 17),
                cell,
                fontsize=10,
                fontname="hebo" if index == 0 else "helv",
            )
            x += width
    y += len(ROWS) * row_height + 34
    page.insert_text((72, y), "Risks", fontsize=13, fontname="hebo")
    y += 24
    for risk in RISKS:
        page.insert_text((80, y), f"- {risk}", fontsize=11)
        y += 17

    page = doc.new_page()
    page_frame(page, 3)
    y = 110
    page.insert_text((72, y), "Outlook", fontsize=17, fontname="hebo")
    y += 28
    for line in OUTLOOK:
        page.insert_text((72, y), line, fontsize=11)
        y += 17

    doc.save(OUT / "annual-report.pdf")
    doc.close()


def quarterly_review_docx() -> None:
    import docx

    d = docx.Document()
    d.core_properties.title = "Quarterly Business Review - Q4 2025"
    d.core_properties.author = "Office of the CFO"

    d.add_heading("Quarterly Business Review", 0)
    d.add_paragraph("Q4 2025 - prepared for the board", style="Subtitle")

    d.add_heading("Executive summary", 1)
    para = d.add_paragraph("The quarter closed at ")
    para.add_run("$4.1M").bold = True
    para.add_run(" in revenue, ")
    para.add_run("18% ahead").bold = True
    para.add_run(" of plan. The beat was concentrated in enterprise, where two expansions closed early. ")
    para.add_run("Margin held").italic = True
    para.add_run(" despite the cost of the new region.")

    d.add_heading("What drove it", 2)
    for text, style in (
        ("Enterprise expansions", "List Bullet"),
        ("Two accounts moved from team to enterprise tier", "List Bullet 2"),
        ("Neither required a discount", "List Bullet 2"),
        ("Lower churn", "List Bullet"),
        ("Gross churn at 1.2%, down from 2.1%", "List Bullet 2"),
        ("Self-serve conversion", "List Bullet"),
    ):
        try:
            d.add_paragraph(text, style=style)
        except KeyError:
            d.add_paragraph(text, style="List Bullet")

    d.add_heading("The numbers", 2)
    table = d.add_table(rows=len(ROWS), cols=4)
    table.style = "Table Grid"
    for r, row in enumerate(ROWS):
        for c, value in enumerate(row):
            table.cell(r, c).text = value

    d.add_heading("Decisions needed", 1)
    for item in (
        "Approve the second engineering pod for the platform team.",
        "Sign off on hedging roughly a fifth of non-dollar billings.",
        "Confirm the pricing change lands before the renewal cohort in March.",
    ):
        d.add_paragraph(item, style="List Number")

    d.add_paragraph("Prepared by the office of the CFO. Figures are unaudited.")
    d.save(OUT / "quarterly-review.docx")


def launch_deck_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.core_properties.title = "Papyrus - launch review"

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Papyrus"
    slide.placeholders[1].text = "Universal document ingestion for agents"
    slide.notes_slide.notes_text_frame.text = (
        "Open with the retrieval failure story - the table that became prose."
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The problem"
    frame = slide.placeholders[1].text_frame
    frame.text = "Models read text. Knowledge lives in files."
    for text, level in (
        ("Text extraction keeps the words", 1),
        ("and throws away the structure", 1),
        ("A flattened table cannot answer a question about the table", 2),
        ("Retrieval quality is capped by ingestion quality", 0),
    ):
        para = frame.add_paragraph()
        para.text = text
        para.level = level
    slide.notes_slide.notes_text_frame.text = "Don't rush this slide. It is the whole pitch."

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Where we are"
    shape = slide.shapes.add_table(5, 3, Inches(0.7), Inches(1.8), Inches(8.6), Inches(3))
    cells = [
        ["Area", "Status", "Note"],
        ["Formats", "22", "PDF through ZIP"],
        ["Tests", "180", "79% coverage"],
        ["Median latency", "180 ms", "48-page PDF"],
        ["LLM calls", "0", "By design"],
    ]
    for r, row in enumerate(cells):
        for c, value in enumerate(row):
            shape.table.cell(r, c).text = value

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "What ships next"
    frame = slide.placeholders[1].text_frame
    frame.text = "OCR for scanned archives"
    for text in ("Layout-aware column detection", "Optional local model for cleanup", "Kubernetes chart"):
        para = frame.add_paragraph()
        para.text = text

    prs.save(OUT / "launch-deck.pptx")


def revenue_model_xlsx() -> None:
    import openpyxl
    from openpyxl.styles import Font

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Revenue"
    ws.append(["Month", "New", "Expansion", "Churn", "Net", "Growth"])
    for cell in ws[1]:
        cell.font = Font(bold=True)

    rows = [
        ("2025-01", 310000, 74000, -22000, 362000, 0.051),
        ("2025-02", 344000, 91000, -18000, 417000, 0.152),
        ("2025-03", 401000, 88000, -25000, 464000, 0.113),
        ("2025-04", 388000, 132000, -19000, 501000, 0.080),
        ("2025-05", 452000, 121000, -31000, 542000, 0.082),
        ("2025-06", 498000, 165000, -27000, 636000, 0.173),
    ]
    for row in rows:
        ws.append(list(row))
    for row in ws.iter_rows(min_row=2, min_col=2, max_col=5):
        for cell in row:
            cell.number_format = '"$"#,##0'
    for row in ws.iter_rows(min_row=2, min_col=6, max_col=6):
        row[0].number_format = "0.0%"
    ws.append([])
    ws.append(["Unaudited management figures."])

    costs = wb.create_sheet("Costs")
    costs.append(["Category", "Monthly", "Share"])
    for cell in costs[1]:
        cell.font = Font(bold=True)
    for row in (
        ("Salaries", 412000, 0.68),
        ("Infrastructure", 61000, 0.10),
        ("Marketing", 88000, 0.15),
        ("Other", 42000, 0.07),
    ):
        costs.append(list(row))
    for row in costs.iter_rows(min_row=2, min_col=2, max_col=2):
        row[0].number_format = '"$"#,##0'
    for row in costs.iter_rows(min_row=2, min_col=3, max_col=3):
        row[0].number_format = "0%"

    scratch = wb.create_sheet("Scratch")
    scratch.append(["working notes", "do not circulate"])
    scratch.sheet_state = "hidden"

    wb.save(OUT / "revenue-model.xlsx")


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    annual_report_pdf()
    quarterly_review_docx()
    launch_deck_pptx()
    revenue_model_xlsx()
    for path in sorted(OUT.iterdir()):
        print(f"  {path.name:26} {path.stat().st_size:>8,} bytes")


if __name__ == "__main__":
    main()
