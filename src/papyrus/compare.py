"""The honest baseline: what you get *without* Papyrus.

Papyrus's whole claim is that structure survives conversion. A claim like
that is only worth anything next to the alternative, so this module
produces the alternative — the naive text extraction people actually write
when they need a document in a prompt:

    with pymupdf.open(path) as pdf:
        text = "".join(page.get_text() for page in pdf)

That is not a straw man. It is one line of every RAG tutorial, it is what
`pdftotext` does, and for Word it is `"\\n".join(p.text for p in
doc.paragraphs)` — which silently omits every table in the document.

Nothing here touches the parsers or the renderer. It reads the same bytes
independently, so the comparison cannot flatter Papyrus by accident.
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field
from typing import Any

from papyrus.converter import ConversionResult
from papyrus.detect import Detection
from papyrus.ir import Document, Table
from papyrus.utils.text import clean, decode

MAX_BASELINE_CHARS = 200_000


@dataclass
class Recovered:
    """What the structured pass found that the naive pass could not."""

    headings: int = 0
    tables: int = 0
    table_cells: int = 0
    lists: int = 0
    list_items: int = 0
    code_blocks: int = 0
    pages: int = 0
    images: int = 0
    running_headers_removed: int = 0
    words: int = 0

    def to_dict(self) -> dict[str, int]:
        return {k: v for k, v in self.__dict__.items()}

    @property
    def headline(self) -> str:
        """One sentence a person would actually repeat."""
        parts: list[str] = []
        if self.tables:
            noun = "table" if self.tables == 1 else "tables"
            parts.append(f"{self.tables} {noun} ({self.table_cells} cells)")
        if self.headings:
            parts.append(f"{self.headings} headings")
        if self.list_items:
            parts.append(f"{self.list_items} list items")
        if self.code_blocks:
            parts.append(f"{self.code_blocks} code blocks")
        if not parts:
            return "Structure preserved, provenance attached."
        recovered = ", ".join(parts[:3])
        suffix = ""
        if self.running_headers_removed:
            suffix = f", and dropped {self.running_headers_removed} repeated page headers"
        return f"Recovered {recovered}{suffix}."


@dataclass
class Comparison:
    filename: str
    format: str
    title: str | None
    baseline: str
    markdown: str
    recovered: Recovered = field(default_factory=Recovered)
    warnings: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        return {
            "filename": self.filename,
            "format": self.format,
            "title": self.title,
            "baseline": self.baseline,
            "markdown": self.markdown,
            "recovered": self.recovered.to_dict(),
            "headline": self.recovered.headline,
            "warnings": self.warnings,
        }


def compare(data: bytes, result: ConversionResult) -> Comparison:
    """Build the side-by-side for an already-converted document."""
    detection = result.detection
    return Comparison(
        filename=result.document.source.filename,
        format=detection.format,
        title=result.document.title,
        baseline=baseline_text(data, detection),
        markdown=result.markdown,
        recovered=measure(result.document),
        warnings=list(result.warnings),
    )


# ── the naive side ───────────────────────────────────────────────────


def baseline_text(data: bytes, detection: Detection) -> str:
    """Extract text the way a one-line script would."""
    try:
        extractor = _EXTRACTORS.get(detection.format, _plain)
        text = extractor(data)
    except Exception:
        # The baseline failing outright is itself a fair result — some
        # formats give a one-liner nothing at all.
        return ""
    return text[:MAX_BASELINE_CHARS].strip()


def _pdf(data: bytes) -> str:
    import pymupdf

    with pymupdf.open(stream=data, filetype="pdf") as pdf:
        return "\n".join(page.get_text() for page in pdf)


def _docx(data: bytes) -> str:
    """Paragraphs only — which is exactly why tables disappear."""
    import docx

    document = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs)


def _pptx(data: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)


def _xlsx(data: bytes) -> str:
    import openpyxl

    book = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    try:
        for name in book.sheetnames:
            for row in book[name].iter_rows(values_only=True):
                cells = [str(v) for v in row if v is not None]
                if cells:
                    lines.append(" ".join(cells))
    finally:
        book.close()
    return "\n".join(lines)


def _html(data: bytes) -> str:
    from bs4 import BeautifulSoup

    return BeautifulSoup(decode(data), "lxml").get_text("\n")


def _plain(data: bytes) -> str:
    return clean(decode(data))


_EXTRACTORS = {
    "pdf": _pdf,
    "docx": _docx,
    "pptx": _pptx,
    "xlsx": _xlsx,
    "html": _html,
    "xml": _html,
    "epub": _plain,
}


# ── the measurement ──────────────────────────────────────────────────


def measure(document: Document) -> Recovered:
    """Count the structure the IR holds that flat text cannot express."""
    found = Recovered(words=document.word_count)

    for block in document.blocks:
        if block.type == "heading":
            found.headings += 1
        elif block.type == "table" and isinstance(block.content, Table):
            found.tables += 1
            table = block.content
            columns = table.width
            found.table_cells += columns * (len(table.rows) + (1 if table.header else 0))
        elif block.type == "list":
            found.lists += 1
            found.list_items += _count_items(block.content)
        elif block.type == "code":
            found.code_blocks += 1
        elif block.type == "page_break":
            found.pages += 1
        elif block.type == "image":
            found.images += 1

    removed = document.metadata.get("removed_running_text")
    if isinstance(removed, list):
        found.running_headers_removed = len(removed) * max(found.pages, 1)
    return found


def _count_items(items: Any) -> int:
    if not isinstance(items, list):
        return 0
    total = 0
    for item in items:
        total += 1
        total += _count_items(getattr(item, "children", []))
    return total
