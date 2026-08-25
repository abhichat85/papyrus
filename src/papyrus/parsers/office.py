"""OOXML: Word, PowerPoint, Excel.

The common trap in all three is document *order*: python-docx hands you
`paragraphs` and `tables` as separate collections, and python-pptx hands
you shapes in z-order rather than reading order. Both are reconstructed
here, because a table that lands at the bottom of the Markdown instead of
between the two paragraphs it belonged to is a silent correctness bug.
"""

from __future__ import annotations

import io
import re
from contextlib import suppress
from typing import Any

from papyrus.config import ConvertOptions
from papyrus.detect import Detection
from papyrus.errors import ParseError
from papyrus.ir import (
    Asset,
    Block,
    Document,
    ListItem,
    heading,
    image,
    list_block,
    page_break,
    paragraph,
    quote,
    table,
)
from papyrus.parsers.base import BaseParser
from papyrus.utils.files import safe_name
from papyrus.utils.tables import alignments, split_header, trim  # noqa: F401
from papyrus.utils.text import clean, slugify

W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"

_HEADING_STYLE = re.compile(r"^(?:Heading|Title|Subtitle|heading)\s*(\d)?", re.IGNORECASE)
_LIST_BULLET = re.compile(r"List\s*(Bullet|Paragraph)", re.IGNORECASE)
_LIST_NUMBER = re.compile(r"List\s*Number", re.IGNORECASE)


# ══════════════════════════════════════════════════════════════════════
# Word
# ══════════════════════════════════════════════════════════════════════


class DocxParser(BaseParser):
    formats = ("docx",)
    label = "Word (.docx)"
    priority = 10

    def parse(self, data: bytes, filename: str, detection: Detection, options: ConvertOptions) -> Document:
        import docx
        from docx.table import Table as DocxTable
        from docx.text.paragraph import Paragraph as DocxParagraph

        doc = self.new_document(data, filename, detection)
        try:
            source = docx.Document(io.BytesIO(data))
        except Exception as exc:
            raise ParseError(f"Could not open Word document: {exc}") from exc

        doc.metadata.update(_docx_core_properties(source))
        if doc.metadata.get("title"):
            doc.title = doc.metadata["title"]

        pending: list[ListItem] = []
        pending_ordered = False

        def flush() -> None:
            nonlocal pending, pending_ordered
            if pending:
                doc.add(list_block(pending, ordered=pending_ordered))
                pending = []

        for item in _iter_body(source):
            if isinstance(item, DocxTable):
                flush()
                block = _docx_table(item)
                if block:
                    doc.add(block)
                continue

            if not isinstance(item, DocxParagraph):
                continue

            text = _docx_paragraph_text(item)
            style = (item.style.name if item.style else "") or ""

            if _has_page_break(item):
                flush()
                doc.add(page_break(len([b for b in doc.blocks if b.type == "page_break"]) + 2))

            if not text.strip():
                continue

            level = _list_level(item)
            if level is not None:
                ordered = bool(_LIST_NUMBER.search(style)) or _is_numbered(item)
                if pending and ordered != pending_ordered:
                    flush()
                pending_ordered = ordered
                _append_nested(pending, ListItem(text), level)
                continue

            flush()
            match = _HEADING_STYLE.match(style)
            if match:
                if style.lower().startswith("title"):
                    doc.title = text
                    continue
                if style.lower().startswith("subtitle"):
                    doc.add(paragraph(f"*{text}*"))
                    continue
                doc.add(heading(text, int(match.group(1) or 1)))
            elif style.lower().startswith(("quote", "intense quote")):
                doc.add(quote(text))
            else:
                doc.add(paragraph(text))

        flush()

        if options.images != "omit":
            _collect_docx_images(source, doc, options)
        if not doc.blocks:
            doc.warn("Word document contained no readable content.")
        return doc


def _iter_body(source: Any) -> list[Any]:
    """Yield paragraphs and tables in true document order."""
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph as DocxParagraph

    body = source.element.body
    out: list[Any] = []
    for child in body.iterchildren():
        if child.tag == f"{W_NS}p":
            out.append(DocxParagraph(child, source))
        elif child.tag == f"{W_NS}tbl":
            out.append(DocxTable(child, source))
    return out


def _docx_paragraph_text(para: Any) -> str:
    """Runs → inline Markdown, with hyperlinks preserved."""
    parts: list[str] = []
    for child in para._p.iterchildren():
        tag = child.tag
        if tag == f"{W_NS}hyperlink":
            label = clean("".join(n.text or "" for n in child.iter(f"{W_NS}t")))
            target = _hyperlink_target(para, child)
            if label:
                parts.append(f"[{label}]({target})" if target else label)
        elif tag == f"{W_NS}r":
            parts.append(_run_text(child))
    text = "".join(parts)
    return clean(text) if text.strip() else ""


def _run_text(run: Any) -> str:
    text = "".join(node.text or "" for node in run.iter(f"{W_NS}t"))
    if not text.strip():
        return text
    props = run.find(f"{W_NS}rPr")
    if props is None:
        return text
    bold = _toggled(props, "b")
    italic = _toggled(props, "i")
    mono = props.find(f"{W_NS}rFonts") is not None and _is_mono(props)
    leading = len(text) - len(text.lstrip())
    trailing = len(text) - len(text.rstrip())
    core = text.strip()
    if mono:
        core = f"`{core}`"
    if bold and italic:
        core = f"***{core}***"
    elif bold:
        core = f"**{core}**"
    elif italic:
        core = f"_{core}_"
    return " " * leading + core + " " * trailing


def _toggled(props: Any, name: str) -> bool:
    node = props.find(f"{W_NS}{name}")
    if node is None:
        return False
    value = node.get(f"{W_NS}val")
    return value not in ("0", "false", "off")


def _is_mono(props: Any) -> bool:
    fonts = props.find(f"{W_NS}rFonts")
    if fonts is None:
        return False
    name = (fonts.get(f"{W_NS}ascii") or "").lower()
    return any(m in name for m in ("courier", "consolas", "mono", "menlo"))


def _hyperlink_target(para: Any, node: Any) -> str:
    rel_id = node.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    if not rel_id:
        return (node.get(f"{W_NS}anchor") and f"#{node.get(f'{W_NS}anchor')}") or ""
    try:
        rel = para.part.rels[rel_id]
        return rel.target_ref if rel.is_external else ""
    except Exception:
        return ""


def _list_level(para: Any) -> int | None:
    """Indent level if this paragraph is a list item, else None."""
    props = para._p.find(f"{W_NS}pPr")
    style = (para.style.name if para.style else "") or ""
    num_pr = props.find(f"{W_NS}numPr") if props is not None else None
    if num_pr is None and not (_LIST_BULLET.search(style) or _LIST_NUMBER.search(style)):
        return None
    if num_pr is not None:
        ilvl = num_pr.find(f"{W_NS}ilvl")
        if ilvl is not None:
            try:
                return max(0, min(6, int(ilvl.get(f"{W_NS}val") or 0)))
            except ValueError:
                return 0
    # Styles carry the depth in their name: "List Bullet 2" is one level in.
    depth = re.search(r"(\d)\s*$", style)
    return max(0, min(6, int(depth.group(1)) - 1)) if depth else 0


def _is_numbered(para: Any) -> bool:
    """True when the list uses a decimal/alpha format rather than a bullet."""
    props = para._p.find(f"{W_NS}pPr")
    if props is None:
        return False
    num_pr = props.find(f"{W_NS}numPr")
    if num_pr is None:
        return False
    num_id_node = num_pr.find(f"{W_NS}numId")
    if num_id_node is None:
        return False
    try:
        num_id = int(num_id_node.get(f"{W_NS}val"))
        numbering = para.part.numbering_part.element
        for num in numbering.iter(f"{W_NS}num"):
            if int(num.get(f"{W_NS}numId")) != num_id:
                continue
            abstract = num.find(f"{W_NS}abstractNumId")
            target = abstract.get(f"{W_NS}val")
            for abs_num in numbering.iter(f"{W_NS}abstractNum"):
                if abs_num.get(f"{W_NS}abstractNumId") != target:
                    continue
                fmt = abs_num.find(f".//{W_NS}numFmt")
                if fmt is not None:
                    return (fmt.get(f"{W_NS}val") or "") != "bullet"
    except Exception:
        return False
    return False


def _append_nested(items: list[ListItem], item: ListItem, level: int) -> None:
    """Place `item` at `level` under the most recent ancestor chain."""
    if level <= 0 or not items:
        items.append(item)
        return
    _append_nested(items[-1].children, item, level - 1)


def _has_page_break(para: Any) -> bool:
    return any(br.get(f"{W_NS}type") == "page" for br in para._p.iter(f"{W_NS}br"))


def _docx_table(source_table: Any) -> Block | None:
    rows: list[list[str]] = []
    for row in source_table.rows:
        cells: list[str] = []
        seen: set[int] = set()
        for cell in row.cells:
            # python-docx returns the *same* cell object for every column a
            # merge spans, so reading them naively repeats the value across
            # the span — a merged total of 190 shows up in both columns as
            # if it were two numbers.
            identity = id(cell._tc)
            if identity in seen:
                cells.append("")
                continue
            seen.add(identity)
            cells.append(clean("\n".join(p.text for p in cell.paragraphs)))
        rows.append(cells)
    rows = trim(rows)
    if not rows:
        return None
    header, body = split_header(rows, _docx_header_hint(source_table, rows))
    return table(body, header=header)


def _docx_header_hint(source_table: Any, rows: list[list[str]]) -> bool | None:
    """Word marks a repeating header row with `w:tblHeader`.

    When the author did not set it, fall back to the convention that a
    multi-row Word table opens with a header — which is true of very nearly
    every table anyone writes in Word — unless row 0 is entirely numeric.
    """
    with suppress(Exception):
        tr_pr = source_table.rows[0]._tr.find(f"{W_NS}trPr")
        if tr_pr is not None and tr_pr.find(f"{W_NS}tblHeader") is not None:
            return True
    if len(rows) < 2:
        return False
    from papyrus.utils.tables import is_numeric

    return not all(is_numeric(c) for c in rows[0] if c.strip())


def _docx_core_properties(source: Any) -> dict[str, Any]:
    out: dict[str, Any] = {}
    try:
        props = source.core_properties
    except Exception:
        return out
    for key in ("title", "author", "subject", "keywords", "category", "comments"):
        value = getattr(props, key, None)
        if value and str(value).strip():
            out["description" if key == "comments" else key] = clean(str(value))
    for key in ("created", "modified"):
        value = getattr(props, key, None)
        if value:
            out[key] = value.isoformat()
    return out


def _collect_docx_images(source: Any, doc: Document, options: ConvertOptions) -> None:
    limits = options.limits
    for rel in source.part.rels.values():
        if "image" not in rel.reltype or rel.is_external:
            continue
        if len(doc.assets) >= limits.max_assets:
            doc.warn(f"Stopped after {limits.max_assets} images.")
            break
        try:
            blob = rel.target_part.blob
        except Exception:
            continue
        if len(blob) < limits.min_asset_bytes:
            continue
        name = safe_name(rel.target_part.partname.split("/")[-1], "image.png")
        asset_id = f"img-{len(doc.assets) + 1:03d}"
        doc.assets.append(Asset(asset_id, f"{asset_id}-{name}", rel.target_part.content_type, blob))
    if doc.assets and options.images in ("extract", "reference"):
        doc.add(heading("Embedded images", 2))
        for asset in doc.assets:
            doc.add(image(asset.asset_id, f"{options.asset_dir}/{asset.filename}", asset.asset_id))


# ══════════════════════════════════════════════════════════════════════
# PowerPoint
# ══════════════════════════════════════════════════════════════════════


class PptxParser(BaseParser):
    formats = ("pptx",)
    label = "PowerPoint (.pptx)"
    priority = 10

    def parse(self, data: bytes, filename: str, detection: Detection, options: ConvertOptions) -> Document:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        doc = self.new_document(data, filename, detection)
        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise ParseError(f"Could not open presentation: {exc}") from exc

        doc.metadata.update(_docx_core_properties(prs))
        if doc.metadata.get("title"):
            doc.title = doc.metadata["title"]

        slides = list(prs.slides)
        limit = options.limits.max_slides
        if len(slides) > limit:
            doc.warn(f"Presentation has {len(slides)} slides; stopped at {limit}.")
            slides = slides[:limit]
        doc.metadata["slides"] = len(slides)

        for index, slide in enumerate(slides, 1):
            title = _slide_title(slide)
            # No label: the heading below already names the slide.
            doc.add(page_break(index))
            doc.add(heading(title or f"Slide {index}", 2, page=index, slide=index))

            for shape in _reading_order(slide, title):
                self._emit_shape(shape, doc, index, options, MSO_SHAPE_TYPE)

            if options.include_speaker_notes:
                notes = _speaker_notes(slide)
                if notes:
                    doc.add(quote(f"**Speaker notes:** {notes}", page=index))

        if not doc.blocks:
            doc.warn("Presentation contained no readable content.")
        return doc

    def _emit_shape(self, shape: Any, doc: Document, index: int, options: ConvertOptions, mso: Any) -> None:
        if shape.has_table:
            block = _pptx_table(shape.table)
            if block:
                block.metadata["page"] = index
                doc.add(block)
            return

        if shape.shape_type == mso.PICTURE and options.images != "omit":
            _emit_pptx_picture(shape, doc, index, options)
            return

        if getattr(shape, "has_chart", False) and shape.has_chart:
            block = _pptx_chart(shape.chart)
            if block:
                block.metadata["page"] = index
                doc.add(block)
            return

        if not shape.has_text_frame:
            return

        bullets: list[ListItem] = []
        for para in shape.text_frame.paragraphs:
            text = clean("".join(run.text for run in para.runs)) or clean(para.text)
            if not text:
                continue
            level = getattr(para, "level", 0) or 0
            _append_nested(bullets, ListItem(text), level)
        if not bullets:
            return
        # A single unindented line is a caption, not a one-item bullet list.
        if len(bullets) == 1 and not bullets[0].children:
            doc.add(paragraph(bullets[0].text, page=index))
        else:
            doc.add(list_block(bullets, page=index))


def _reading_order(slide: Any, title: str | None) -> list[Any]:
    """Shapes top-to-bottom, left-to-right — not the z-order pptx gives us."""
    shapes = []
    for shape in slide.shapes:
        if title and shape == slide.shapes.title:
            continue
        top = shape.top if shape.top is not None else 0
        left = shape.left if shape.left is not None else 0
        shapes.append((top, left, shape))
    shapes.sort(key=lambda t: (t[0], t[1]))
    return [s for _, _, s in shapes]


def _slide_title(slide: Any) -> str | None:
    with suppress(Exception):
        if slide.shapes.title and slide.shapes.title.text.strip():
            return clean(slide.shapes.title.text)
    return None


def _speaker_notes(slide: Any) -> str:
    with suppress(Exception):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            return clean(slide.notes_slide.notes_text_frame.text)
    return ""


def _pptx_table(source_table: Any) -> Block | None:
    rows = [[clean(cell.text) for cell in row.cells] for row in source_table.rows]
    rows = trim(rows)
    if not rows:
        return None
    hint = bool(getattr(source_table, "first_row", False)) or None
    header, body = split_header(rows, hint)
    return table(body, header=header)


def _pptx_chart(chart: Any) -> Block | None:
    """Charts carry their source data — recover it as a table."""
    try:
        categories = [str(c) for c in chart.plots[0].categories]
        series = list(chart.series)
    except Exception:
        return None
    if not categories or not series:
        return None
    header = ["Category"] + [str(s.name or f"Series {i}") for i, s in enumerate(series, 1)]
    rows: list[list[str]] = []
    for i, category in enumerate(categories):
        row = [category]
        for s in series:
            values = list(s.values)
            row.append("" if i >= len(values) or values[i] is None else str(values[i]))
        rows.append(row)
    return table(rows, header=header, caption=f"Chart data: {chart.chart_type}")


def _emit_pptx_picture(shape: Any, doc: Document, index: int, options: ConvertOptions) -> None:
    limits = options.limits
    if len(doc.assets) >= limits.max_assets:
        return
    try:
        blob = shape.image.blob
        ext = shape.image.ext
        content_type = shape.image.content_type
    except Exception:
        return
    if len(blob) < limits.min_asset_bytes:
        return
    asset_id = f"img-{len(doc.assets) + 1:03d}"
    filename = f"{asset_id}-slide{index}.{ext}"
    doc.assets.append(Asset(asset_id, filename, content_type, blob))
    alt = clean(getattr(shape, "name", "") or "") or f"Slide {index} image"
    doc.add(image(alt, f"{options.asset_dir}/{filename}", asset_id, page=index))


# ══════════════════════════════════════════════════════════════════════
# Excel
# ══════════════════════════════════════════════════════════════════════


class XlsxParser(BaseParser):
    formats = ("xlsx",)
    label = "Excel (.xlsx)"
    priority = 10

    def parse(self, data: bytes, filename: str, detection: Detection, options: ConvertOptions) -> Document:
        import openpyxl

        doc = self.new_document(data, filename, detection)
        try:
            book = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True, keep_links=False)
        except Exception as exc:
            raise ParseError(f"Could not open workbook: {exc}") from exc

        limits = options.limits
        sheet_names = list(book.sheetnames)
        doc.metadata["sheets"] = sheet_names
        budget = limits.max_table_cells
        blank_cells = False

        try:
            for name in sheet_names:
                sheet = book[name]
                hidden = getattr(sheet, "sheet_state", "visible") != "visible"
                if hidden and not options.include_hidden_sheets:
                    doc.warn(f"Skipped hidden sheet '{name}'.")
                    continue

                rows, truncated = _sheet_rows(sheet, limits, budget)
                budget -= sum(len(r) for r in rows)
                if any(not cell.strip() for row in rows for cell in row):
                    blank_cells = True
                rows = trim(rows)
                if not rows:
                    continue

                doc.add(heading(f"Sheet: {name}", 2, sheet=name, anchor=slugify(name)))
                if truncated:
                    doc.warn(f"Sheet '{name}' truncated at {limits.max_sheet_rows} rows.")

                for block in _split_sheet_regions(rows):
                    doc.add(block)
        finally:
            book.close()

        # An empty cell might be an empty cell, or it might be a formula
        # whose result was never cached — a workbook written by a script has
        # no cached values at all, so whole columns come back blank. Saying
        # which it is turns a confusing silence into a fixable instruction.
        if blank_cells:
            with_formulas = _sheets_with_formulas(data)
            if with_formulas:
                listed = ", ".join(f"'{n}'" for n in with_formulas[:4])
                doc.warn(
                    f"Formula results are missing from {listed}. Excel stores the last "
                    "computed value, and this workbook has none — open and save it in a "
                    "spreadsheet app, then convert again."
                )

        if not doc.blocks:
            doc.warn("Workbook contained no data.")
        return doc


def _sheets_with_formulas(data: bytes) -> list[str]:
    """Sheets containing formula cells, read without the cached-value pass.

    Only called when at least one cell came back blank, so the second load
    is paid for exactly when it can explain something.
    """
    import openpyxl

    found: list[str] = []
    try:
        book = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=False)
    except Exception:
        return found
    try:
        for name in book.sheetnames:
            sheet = book[name]
            try:
                for row in sheet.iter_rows(values_only=True):
                    if any(isinstance(v, str) and v.startswith("=") for v in row):
                        found.append(name)
                        break
            except Exception:
                continue
    finally:
        book.close()
    return found


def _sheet_rows(sheet: Any, limits: Any, budget: int) -> tuple[list[list[str]], bool]:
    rows: list[list[str]] = []
    truncated = False
    for row in sheet.iter_rows():
        if len(rows) >= limits.max_sheet_rows or budget <= 0:
            truncated = True
            break
        values = [_cell(c) for c in row[: limits.max_sheet_cols]]
        budget -= len(values)
        rows.append(values)
    return rows, truncated


def _cell(cell: Any) -> str:
    """Render a cell the way the spreadsheet displays it.

    The stored value is not what the sheet shows: a percentage is stored as
    0.05, currency as a bare float. Dropping the number format turns "5%"
    into "0.05" — the kind of quiet corruption that is worse than a missing
    cell, because nothing downstream can tell it happened.
    """
    value = getattr(cell, "value", cell)
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"

    fmt = (getattr(cell, "number_format", "") or "").lower()
    if isinstance(value, int | float) and not isinstance(value, bool):
        if "%" in fmt:
            scaled = value * 100
            decimals = _format_decimals(fmt)
            return f"{scaled:.{decimals}f}%".replace(".00%", "%") if decimals else f"{scaled:g}%"
        symbol = _currency_symbol(fmt)
        magnitude = abs(value)
        number = str(int(magnitude)) if float(magnitude).is_integer() else f"{magnitude:g}"
        sign = "-" if value < 0 else ""
        # The minus belongs outside the symbol: -$22,000, never $-22,000.
        return f"{sign}{symbol}{number}" if symbol else f"{sign}{number}"

    if hasattr(value, "isoformat"):
        text = value.isoformat()
        # Excel stores dates as datetimes; drop a midnight time component
        # the sheet is not displaying.
        return text[:10] if text.endswith("T00:00:00") and "h" not in fmt else text
    return clean(str(value))


def _format_decimals(fmt: str) -> int:
    if "." not in fmt:
        return 0
    tail = fmt.split(".", 1)[1]
    return len(tail) - len(tail.lstrip("0#"))


def _currency_symbol(fmt: str) -> str:
    return next((s for s in ("$", "\u20ac", "\u00a3", "\u00a5", "\u20b9") if s in fmt), "")


def _sheet_header_hint(region: list[list[str]]) -> bool | None:
    """A spreadsheet region opens with a header when its first row is fully
    populated and, unlike the rows beneath it, carries no numbers."""
    if len(region) < 2:
        return False
    first = region[0]
    if not all(c.strip() for c in first):
        return None
    from papyrus.utils.tables import is_numeric

    if any(is_numeric(c) for c in first):
        return None
    return True


def _split_sheet_regions(rows: list[list[str]]) -> list[Block]:
    """Split a sheet on blank-row gutters.

    People put three unrelated tables on one sheet separated by empty rows.
    Rendering that as a single ragged table produces Markdown nobody — human
    or model — can read.
    """
    regions: list[list[list[str]]] = []
    current: list[list[str]] = []
    for row in rows:
        if any(c.strip() for c in row):
            current.append(row)
        elif current:
            regions.append(current)
            current = []
    if current:
        regions.append(current)

    blocks: list[Block] = []
    for region in regions:
        region = trim(region)
        if not region:
            continue
        # A lone one-column region is prose (a note or title), not a table.
        if len(region[0]) == 1:
            for row in region:
                if row[0].strip():
                    blocks.append(paragraph(row[0]))
            continue
        header, body = split_header(region, _sheet_header_hint(region))
        blocks.append(table(body, header=header))
    return blocks
