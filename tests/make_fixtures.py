"""Generate binary test fixtures.

Binary fixtures are built, not committed: a .docx in git is an opaque blob
nobody can review, and regenerating them keeps the expectations in the
tests readable. Run `python tests/make_fixtures.py` after changing this.
"""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

FIXTURES = Path(__file__).parent / "fixtures"


def make_docx() -> None:
    import docx

    d = docx.Document()
    d.core_properties.title = "Quarterly Review"
    d.core_properties.author = "Finance Team"
    d.add_heading("Quarterly Review", 0)
    d.add_heading("Executive Summary", 1)
    p = d.add_paragraph("Revenue grew ")
    p.add_run("41%").bold = True
    p.add_run(" year over year, ahead of the ")
    p.add_run("conservative").italic = True
    p.add_run(" plan.")
    d.add_heading("Highlights", 2)
    for text, style in (
        ("Enterprise pipeline doubled", "List Bullet"),
        ("Driven by two expansions", "List Bullet 2"),
        ("Churn fell to 1.2%", "List Bullet"),
        ("Two new regions", "List Bullet"),
    ):
        d.add_paragraph(text, style=style)
    d.add_heading("Numbers", 2)
    table = d.add_table(rows=3, cols=3)
    data = [["Metric", "2024", "2025"], ["Revenue", "10000000", "14100000"], ["EBITDA", "900000", "2100000"]]
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    d.add_paragraph("Prepared by the finance team.")
    d.save(FIXTURES / "sample.docx")


def make_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Papyrus Launch"
    slide.placeholders[1].text = "Universal document ingestion"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Why now"
    body = slide.placeholders[1].text_frame
    body.text = "Agents need clean text"
    for text, level in (
        ("PDFs are the worst offender", 1),
        ("Tables carry the meaning", 1),
        ("Provenance matters", 0),
    ):
        para = body.add_paragraph()
        para.text = text
        para.level = level
    slide.notes_slide.notes_text_frame.text = "Open with the retrieval failure story."

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Numbers"
    shape = slide.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(6), Inches(2))
    cells = [["Metric", "Value"], ["Formats", "15"], ["Latency", "180ms"]]
    for r, row in enumerate(cells):
        for c, value in enumerate(row):
            shape.table.cell(r, c).text = value
    prs.save(FIXTURES / "sample.pptx")


def make_xlsx() -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Revenue"
    for row in [
        ["Month", "Revenue", "Growth"],
        ["Jan", 100000, 0.05],
        ["Feb", -120000, 0.2],  # negative, to pin the currency sign down
        ["Mar", 140000, 0.17],
    ]:
        ws.append(row)
    for row in ws.iter_rows(min_row=2, min_col=3, max_col=3):
        row[0].number_format = "0%"
    for row in ws.iter_rows(min_row=2, min_col=2, max_col=2):
        row[0].number_format = '"$"#,##0'
    ws.append([])
    ws.append(["Note: unaudited"])

    ws2 = wb.create_sheet("Expenses")
    for row in [["Category", "Amount"], ["Salaries", 50000], ["Marketing", 12000]]:
        ws2.append(row)

    hidden = wb.create_sheet("Scratch")
    hidden.append(["do not read", 1])
    hidden.sheet_state = "hidden"
    wb.save(FIXTURES / "sample.xlsx")


def make_pdf() -> None:
    import pymupdf

    doc = pymupdf.open()
    doc.set_metadata({"title": "Papyrus Test Report", "author": "Einstein Labs"})

    for page_no in (1, 2):
        page = doc.new_page()
        page.insert_text((72, 50), "Papyrus Test Report | Confidential", fontsize=8)
        y = 100
        if page_no == 1:
            page.insert_text((72, y), "Annual Report", fontsize=24, fontname="hebo")
            y += 44
            page.insert_text((72, y), "Executive Summary", fontsize=16, fontname="hebo")
            y += 30
            for line in (
                "Revenue grew forty one percent year over year, ahead of the",
                "conservative plan set at the start of the period.",
            ):
                page.insert_text((72, y), line, fontsize=11)
                y += 16
            y += 12
            for bullet in (
                "- Enterprise pipeline doubled",
                "- Churn fell to 1.2 percent",
                "- Two new regions opened",
            ):
                page.insert_text((72, y), bullet, fontsize=11)
                y += 16
        else:
            page.insert_text((72, y), "Outlook", fontsize=16, fontname="hebo")
            y += 30
            page.insert_text((72, y), "The coming year is expected to compound on this base.", fontsize=11)
        page.insert_text((300, 780), f"{page_no}", fontsize=8)

    # A third page where a table sits BETWEEN two headed sections, to pin
    # down reading order: text and tables come from separate extraction
    # passes and must still interleave correctly.
    page = doc.new_page()
    page.insert_text((72, 50), "Papyrus Test Report | Confidential", fontsize=8)
    page.insert_text((72, 100), "Financials", fontsize=16, fontname="hebo")
    top = 130
    for r, row in enumerate([["Metric", "2025"], ["Revenue", "14100000"], ["EBITDA", "2100000"]]):
        for c, cell in enumerate(row):
            rect = pymupdf.Rect(72 + c * 120, top + r * 24, 72 + (c + 1) * 120, top + (r + 1) * 24)
            page.draw_rect(rect, color=(0.7, 0.7, 0.7), width=0.6)
            page.insert_text((78 + c * 120, top + r * 24 + 16), cell, fontsize=10)
    page.insert_text((72, top + 3 * 24 + 40), "Notes", fontsize=16, fontname="hebo")
    page.insert_text((72, top + 3 * 24 + 66), "Figures are unaudited.", fontsize=11)
    page.insert_text((300, 780), "3", fontsize=8)

    doc.save(FIXTURES / "sample.pdf")
    doc.close()


def make_text_fixtures() -> None:
    (FIXTURES / "sample.csv").write_text(
        "Month,Revenue,Growth\nJan,100000,5%\nFeb,120000,20%\nMar,140000,17%\n", encoding="utf-8"
    )
    (FIXTURES / "sample.tsv").write_text("a\tb\n1\t2\n3\t4\n", encoding="utf-8")
    (FIXTURES / "sample.txt").write_text(
        "PROJECT BRIEF\n\nThe goal is a universal ingestion engine.\nIt must never crash on bad input.\n\n"
        "- deterministic\n- local first\n- fast\n",
        encoding="utf-8",
    )
    (FIXTURES / "sample.md").write_text("# Title\n\nBody **text** here.\n\n- one\n- two\n", encoding="utf-8")
    (FIXTURES / "sample.json").write_text(
        json.dumps(
            {
                "product": "Papyrus",
                "version": "0.1.0",
                "formats": ["pdf", "docx", "pptx"],
                "benchmarks": [
                    {"file": "report.pdf", "ms": 180, "pages": 48},
                    {"file": "deck.pptx", "ms": 40, "pages": 22},
                ],
            },
            indent=2,
        ),
        encoding="utf-8",
    )
    (FIXTURES / "sample.jsonl").write_text(
        '{"id":1,"name":"a"}\n{"id":2,"name":"b"}\n{"id":3,"name":"c"}\n', encoding="utf-8"
    )
    (FIXTURES / "sample.html").write_text(
        "<!doctype html><html><head><title>Doc Title</title>"
        '<meta name="description" content="A test page">'
        "</head><body><nav>skip me</nav><main>"
        "<h1>Doc Title</h1><p>Intro with <a href='https://example.com'>a link</a> and <strong>bold</strong>.</p>"
        "<h2>List</h2><ul><li>one<ul><li>nested</li></ul></li><li>two</li></ul>"
        "<table><thead><tr><th>K</th><th>V</th></tr></thead><tbody><tr><td>a</td><td>1</td></tr></tbody></table>"
        "<pre><code class='language-python'>print('hi')</code></pre>"
        "</main><footer>skip me too</footer></body></html>",
        encoding="utf-8",
    )
    (FIXTURES / "sample.py").write_text('def hello():\n    return "world"\n', encoding="utf-8")
    (FIXTURES / "sample.eml").write_text(
        "From: alice@example.com\nTo: bob@example.com\nSubject: Launch plan\n"
        "Date: Mon, 1 Jan 2026 09:00:00 +0000\nContent-Type: text/plain; charset=utf-8\n\n"
        "Hi Bob,\n\nThe plan is attached.\n\n- ship Monday\n- announce Tuesday\n\nAlice\n",
        encoding="utf-8",
    )
    (FIXTURES / "sample.rtf").write_text(
        r"{\rtf1\ansi{\fonttbl\f0\fswiss Helvetica;}\f0\pard Hello \b bold\b0  world.\par Second paragraph.\par}",
        encoding="utf-8",
    )
    (FIXTURES / "sample.ipynb").write_text(
        json.dumps(
            {
                "cells": [
                    {
                        "cell_type": "markdown",
                        "metadata": {},
                        "source": ["# Notebook\n", "\n", "Some prose.\n"],
                    },
                    {
                        "cell_type": "code",
                        "execution_count": 1,
                        "metadata": {},
                        "source": ["print('hello')\n"],
                        "outputs": [{"output_type": "stream", "name": "stdout", "text": ["hello\n"]}],
                    },
                ],
                "metadata": {"kernelspec": {"language": "python", "display_name": "Python 3"}},
                "nbformat": 4,
                "nbformat_minor": 5,
            },
            indent=1,
        ),
        encoding="utf-8",
    )
    with zipfile.ZipFile(FIXTURES / "sample.zip", "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("notes.md", "# Notes\n\nInside the archive.\n")
        z.writestr("data.csv", "a,b\n1,2\n")
        z.writestr("__MACOSX/junk", "ignore me")


def main() -> None:
    FIXTURES.mkdir(parents=True, exist_ok=True)
    make_text_fixtures()
    make_docx()
    make_pptx()
    make_xlsx()
    make_pdf()
    print(f"fixtures written to {FIXTURES}")


if __name__ == "__main__":
    main()
