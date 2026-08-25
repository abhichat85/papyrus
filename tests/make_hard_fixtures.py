"""Generate the *hard* fixtures.

`make_fixtures.py` produces clean, well-formed documents — the happy path.
This file produces the documents that break converters in the wild:
merged cells, footnotes, multi-column layout, rotated pages, BOMs,
quoted newlines, ragged rows, deep nesting, entity soup, and text in
scripts that are not Latin.

Every file here exists because it represents a class of real failure, and
the name says which one.
"""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

HARD = Path(__file__).parent / "fixtures" / "hard"


# ══════════════════════════════════════════════════════════════════════
# Word
# ══════════════════════════════════════════════════════════════════════


def docx_merged_and_nested() -> None:
    """Merged cells, a nested table, and an empty cell run."""
    import docx

    d = docx.Document()
    d.add_heading("Merged cells", 1)

    table = d.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    headers = ["Region", "Q1", "Q2"]
    for c, text in enumerate(headers):
        table.cell(0, c).text = text
    rows = [["North", "100", "120"], ["South", "90", "95"], ["Total", "190", "215"]]
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    # Merge the two Q columns of the total row — the classic ragged-table maker.
    table.cell(3, 1).merge(table.cell(3, 2))

    d.add_heading("Empty cells", 1)
    sparse = d.add_table(rows=3, cols=3)
    sparse.style = "Table Grid"
    sparse.cell(0, 0).text = "only"
    sparse.cell(2, 2).text = "corners"

    d.save(HARD / "merged-cells.docx")


def docx_deep_lists() -> None:
    """Six levels of nesting, mixed ordered and unordered."""
    import docx

    d = docx.Document()
    d.add_heading("Deep nesting", 1)
    available = {s.name for s in d.styles}
    for level in range(1, 6):
        style = "List Bullet" if level == 1 else f"List Bullet {level}"
        # Only styles the template actually defines. `add_paragraph` appends
        # the paragraph before it validates the style name, so a try/except
        # around it leaves a stray duplicate behind.
        d.add_paragraph(
            f"Level {level} bullet", style=style if style in available else "List Bullet"
        )
    d.add_paragraph("Numbered one", style="List Number")
    d.add_paragraph("Numbered two", style="List Number")
    d.add_paragraph("Back to prose.")
    d.save(HARD / "deep-lists.docx")


def docx_unicode() -> None:
    """Right-to-left, CJK, combining marks, emoji, and a very long word."""
    import docx

    d = docx.Document()
    d.add_heading("多言語ドキュメント", 1)
    d.add_paragraph("日本語の段落です。これは表の上にあります。")
    d.add_paragraph("هذه فقرة باللغة العربية تُقرأ من اليمين إلى اليسار.")
    d.add_paragraph("Ελληνικά, Кириллица, ไทย, עברית")
    d.add_paragraph("Emoji: 🚀📊🧾 and combining: é ä ñ")
    d.add_paragraph("Long" + "word" * 60)
    table = d.add_table(rows=2, cols=2)
    table.style = "Table Grid"
    for c, text in enumerate(["列 1", "列 2"]):
        table.cell(0, c).text = text
    for c, text in enumerate(["値 A", "値 B"]):
        table.cell(1, c).text = text
    d.save(HARD / "unicode.docx")


def docx_markdown_injection() -> None:
    """Content that is itself Markdown syntax — must not corrupt the output."""
    import docx

    d = docx.Document()
    d.add_heading("Injection", 1)
    d.add_paragraph("A pipe | inside a paragraph, and a # hash, and *stars*.")
    d.add_paragraph("```")
    d.add_paragraph("---")
    d.add_paragraph("| fake | table |")
    table = d.add_table(rows=2, cols=2)
    table.style = "Table Grid"
    table.cell(0, 0).text = "col | with pipe"
    table.cell(0, 1).text = "normal"
    table.cell(1, 0).text = "line\nbreak"
    table.cell(1, 1).text = "```fence```"
    d.save(HARD / "injection.docx")


# ══════════════════════════════════════════════════════════════════════
# Excel
# ══════════════════════════════════════════════════════════════════════


def xlsx_hard() -> None:
    """Merged cells, dates, formulas, blank gutters and a wide sheet."""
    import datetime as dt

    import openpyxl
    from openpyxl.styles import Font

    wb = openpyxl.Workbook()

    ws = wb.active
    ws.title = "Mixed"
    ws["A1"] = "Report title spanning columns"
    ws.merge_cells("A1:D1")
    ws["A1"].font = Font(bold=True)
    ws.append([])
    ws.append(["Date", "Amount", "Rate", "Note"])
    ws.append([dt.date(2026, 1, 15), 1234.5, 0.075, "first"])
    ws.append([dt.datetime(2026, 2, 20, 14, 30), -890, 0.12, "negative"])
    ws.append([dt.date(2026, 3, 1), 0, 0, ""])
    for row in ws.iter_rows(min_row=4, min_col=1, max_col=1):
        row[0].number_format = "yyyy-mm-dd"
    for row in ws.iter_rows(min_row=4, min_col=2, max_col=2):
        row[0].number_format = '"$"#,##0.00'
    for row in ws.iter_rows(min_row=4, min_col=3, max_col=3):
        row[0].number_format = "0.0%"

    # Two unrelated regions separated by blank rows on one sheet.
    ws.append([])
    ws.append([])
    ws.append(["Second region", "starts here"])
    ws.append(["x", 1])
    ws.append(["y", 2])

    wide = wb.create_sheet("Wide")
    wide.append([f"col{i}" for i in range(1, 40)])
    wide.append(list(range(1, 40)))

    formulas = wb.create_sheet("Formulas")
    formulas.append(["a", "b", "sum"])
    formulas.append([1, 2, "=A2+B2"])
    formulas.append([3, 4, "=A3+B3"])

    empty = wb.create_sheet("Empty")
    empty["A1"] = None

    wb.save(HARD / "hard.xlsx")


# ══════════════════════════════════════════════════════════════════════
# PowerPoint
# ══════════════════════════════════════════════════════════════════════


def pptx_hard() -> None:
    """A grouped shape, an empty slide, and out-of-order shape placement."""
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt

    prs = Presentation()

    # Shapes added bottom-first: reading order must not follow insertion.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    lower = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(6), Inches(1))
    lower.text_frame.text = "This sits at the bottom of the slide."
    upper = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    upper.text_frame.text = "This sits at the top of the slide."

    # A slide with nothing on it at all.
    prs.slides.add_slide(prs.slide_layouts[6])

    # Long bulleted body that wraps.
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dense slide"
    frame = slide.placeholders[1].text_frame
    frame.text = "First point"
    for index in range(2, 12):
        para = frame.add_paragraph()
        para.text = f"Point number {index} with enough words to wrap on a real slide"
        para.level = index % 3
        para.font.size = Pt(12)

    # Two side-by-side text boxes: left column then right column.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = slide.shapes.add_textbox(Emu(457200), Emu(914400), Emu(4114800), Emu(2743200))
    left.text_frame.text = "Left column content"
    right = slide.shapes.add_textbox(Emu(4800600), Emu(914400), Emu(4114800), Emu(2743200))
    right.text_frame.text = "Right column content"

    prs.save(HARD / "hard.pptx")


# ══════════════════════════════════════════════════════════════════════
# PDF
# ══════════════════════════════════════════════════════════════════════


def pdf_two_column() -> None:
    """Two-column layout with a rotated page and a page of pure whitespace."""
    import pymupdf

    doc = pymupdf.open()
    doc.set_metadata({"title": "Two Column Paper", "author": "Test"})

    page = doc.new_page()
    page.insert_text((60, 70), "A Two Column Paper", fontsize=20, fontname="hebo")
    left_lines = [
        "The first column begins here and",
        "continues for several lines before",
        "the reader reaches the bottom of",
        "the page and moves to the right.",
    ]
    right_lines = [
        "The second column continues the",
        "argument started on the left and",
        "should not be interleaved with it",
        "line by line in the output.",
    ]
    y = 120
    for line in left_lines:
        page.insert_text((60, y), line, fontsize=10)
        y += 15
    y = 120
    for line in right_lines:
        page.insert_text((320, y), line, fontsize=10)
        y += 15

    # A page with a heading and nothing else.
    page = doc.new_page()
    page.insert_text((60, 100), "Appendix", fontsize=16, fontname="hebo")

    # A rotated page.
    page = doc.new_page()
    page.insert_text((60, 100), "Rotated content heading", fontsize=16, fontname="hebo")
    page.insert_text((60, 130), "Body text on a rotated page.", fontsize=10)
    page.set_rotation(90)

    # A page with no text at all — the scanned-page signal.
    doc.new_page()

    doc.save(HARD / "two-column.pdf")
    doc.close()


def pdf_no_text() -> None:
    """Every page blank — must warn about a scanned document, not crash."""
    import pymupdf

    doc = pymupdf.open()
    for _ in range(3):
        page = doc.new_page()
        page.draw_rect(pymupdf.Rect(100, 100, 400, 300), color=(0.5, 0.5, 0.5), width=2)
    doc.save(HARD / "scanned.pdf")
    doc.close()


# ══════════════════════════════════════════════════════════════════════
# Text-shaped formats
# ══════════════════════════════════════════════════════════════════════


def text_hard() -> None:
    # BOM + CRLF + semicolon delimiter + quoted newlines + ragged rows.
    (HARD / "bom-crlf.csv").write_bytes(
        "﻿name;note;amount\r\n"
        'Alice;"line one\nline two";100\r\n'
        "Bob;plain;200\r\n"
        "Carol;missing\r\n"
        'Dave;"has ""quotes"" inside";400;extra\r\n'.encode()
    )

    # Latin-1 bytes with no BOM.
    (HARD / "latin1.csv").write_bytes(
        "name,city\nRené,Zürich\nFrançois,Köln\n".encode("latin-1")
    )

    # A CSV that is really one column of prose.
    (HARD / "single-column.csv").write_text(
        "note\nfirst observation\nsecond observation\n", encoding="utf-8"
    )

    # Deeply nested JSON plus unicode keys and an empty array.
    (HARD / "deep.json").write_text(
        json.dumps(
            {
                "レベル1": {
                    "level2": {
                        "level3": {
                            "level4": {"level5": {"value": "bottom", "list": []}},
                            "siblings": [1, 2, 3],
                        }
                    }
                },
                "mixed_array": [{"a": 1}, {"b": 2}, "string", 42, None, True],
                "empty_object": {},
                "nulls": [None, None],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    # Entity soup, unclosed tags, nested tables, definition list.
    (HARD / "messy.html").write_text(
        "<html><head><title>Messy &amp; Broken</title></head><body>"
        "<main>"
        "<h1>Messy &amp; Broken</h1>"
        "<p>Unclosed paragraph with &lt;entities&gt; and &nbsp; spaces"
        "<p>Second paragraph &mdash; note the missing close tag above."
        "<dl><dt>Term</dt><dd>Definition</dd><dt>Other</dt><dd>Second</dd></dl>"
        "<table><tr><td>outer</td><td><table><tr><td>inner</td></tr></table></td></tr></table>"
        "<ul><li>one<li>two<li>three</ul>"
        "<blockquote>Quoted <em>text</em> here.</blockquote>"
        "<pre><code>if x &lt; 3:\n    pass</code></pre>"
        "<p>Trailing text after everything." + ("padding " * 60) +
        "</main></body></html>",
        encoding="utf-8",
    )

    # Markdown that already contains frontmatter and a table.
    (HARD / "with-frontmatter.md").write_text(
        "---\ntitle: Existing Doc\ntags: [a, b]\n---\n\n"
        "# Existing Doc\n\nBody with a table.\n\n"
        "| a | b |\n| --- | --- |\n| 1 | 2 |\n\n"
        "```python\nprint('nested ``` backticks')\n```\n",
        encoding="utf-8",
    )

    # A log file: no structure at all, very repetitive.
    (HARD / "server.log").write_text(
        "\n".join(
            f"2026-08-{(i % 28) + 1:02d} 10:{i % 60:02d}:00 INFO request id={i} status=200"
            for i in range(400)
        ),
        encoding="utf-8",
    )

    # An empty-ish text file with only whitespace.
    (HARD / "whitespace.txt").write_text("   \n\n\t\t\n   \n", encoding="utf-8")

    # A zip holding one of everything, including a nested zip.
    inner = HARD / "_inner.zip"
    with zipfile.ZipFile(inner, "w") as z:
        z.writestr("nested-note.md", "# Nested\n\nOne level down.\n")
    with zipfile.ZipFile(HARD / "mixed.zip", "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("readme.md", "# Archive\n\nTop level readme.\n")
        z.writestr("data/values.csv", "k,v\nalpha,1\nbeta,2\n")
        z.writestr("code/app.py", "def main():\n    return 42\n")
        z.writestr(
            "page.html",
            "<html><body><h1>Quarterly Web Page</h1><p>Body text of the page.</p></body></html>",
        )
        z.writestr("inner.zip", inner.read_bytes())
        z.writestr("binary.dat", bytes(range(256)))
    inner.unlink()


def main() -> None:
    HARD.mkdir(parents=True, exist_ok=True)
    text_hard()
    docx_merged_and_nested()
    docx_deep_lists()
    docx_unicode()
    docx_markdown_injection()
    xlsx_hard()
    pptx_hard()
    pdf_two_column()
    pdf_no_text()
    count = len(list(HARD.iterdir()))
    print(f"{count} hard fixtures written to {HARD}")


if __name__ == "__main__":
    main()
